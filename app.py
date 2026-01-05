import os
import logging
import datetime
import uuid
import time
import threading
import json
from io import BytesIO

# --- FLASK IMPORTS ---
from flask import Flask, render_template, request, jsonify, Response, session, redirect, url_for, send_from_directory
from flask_sqlalchemy import SQLAlchemy
from werkzeug.security import generate_password_hash, check_password_hash
from dotenv import load_dotenv

# --- SYSTEM & UTILS ---
import psutil

# --- AI & MEDIA IMPORTS ---
from groq import Groq
from gtts import gTTS
from xhtml2pdf import pisa
from pptx import Presentation
import speech_recognition as sr
import PIL.Image
from moviepy.video.io.VideoFileClip import VideoFileClip

# --- SETUP ENV ---
basedir = os.path.abspath(os.path.dirname(__file__))
load_dotenv(os.path.join(basedir, '.env'))

app = Flask(__name__)
app.secret_key = os.environ.get("SECRET_KEY", "super_secret_key_123")

# --- VERCEL SPECIFIC CONFIGURATION ---
# Check if running on Vercel
IS_VERCEL = os.environ.get('VERCEL') == '1'

if IS_VERCEL:
    # Vercel only allows writing to /tmp
    STATIC_FOLDER = '/tmp'
    # Use a temporary database in /tmp (WARNING: Data resets on redeploy)
    app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:////tmp/users.db'
else:
    # Local development settings
    STATIC_FOLDER = os.path.join(basedir, 'static')
    if not os.path.exists(STATIC_FOLDER):
        os.makedirs(STATIC_FOLDER)
    app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///' + os.path.join(basedir, 'users.db')

app.config['SQLALCHEMY_TRACK_MODIFICATIONS'] = False

# FIX: Allow uploads up to 100MB (Critical for Video to Audio)
app.config['MAX_CONTENT_LENGTH'] = 100 * 1024 * 1024 

db = SQLAlchemy(app)

# --- AI CONFIG ---
API_KEY = os.environ.get("GROQ_API_KEY")

# --- ADMIN CONFIG ---
ADMIN_USER = os.environ.get("ADMIN_USER", "admin")
ADMIN_PASS = os.environ.get("ADMIN_PASS", "admin123") 

# ==============================================================================
#                               DATABASE MODELS
# ==============================================================================

class User(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(80), unique=True, nullable=False)
    password_hash = db.Column(db.String(120), nullable=False)
    # Relationship to access user's activities easily
    activities = db.relationship('ActivityLog', backref='user', lazy=True)

class ActivityLog(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=True) # Nullable for Admin/Guest actions
    activity_type = db.Column(db.String(50), nullable=False)  # e.g., 'text_to_pdf', 'chat'
    details = db.Column(db.String(200)) # e.g., 'Converted report.docx'
    timestamp = db.Column(db.DateTime, default=datetime.datetime.utcnow)

# Create Database Tables
# We do this inside a try/except block to ensure it works on Vercel startup
with app.app_context():
    try:
        db.create_all()
    except Exception as e:
        print(f"DB Init Error: {e}")

# --- HELPER: LOGGING TO DB ---
def log_activity(activity_type, details=""):
    try:
        # 1. Update Global Stats (In-Memory for Live Dashboard)
        if activity_type in global_stats:
            global_stats[activity_type] += 1
            
        # 2. Log to Database (Permanent Record)
        current_username = session.get('user_name')
        user_id = None
        
        if current_username and current_username != "Administrator":
            user = User.query.filter_by(username=current_username).first()
            if user:
                user_id = user.id
        
        # Create log entry
        new_log = ActivityLog(
            user_id=user_id,
            activity_type=activity_type,
            details=details,
            timestamp=datetime.datetime.now()
        )
        db.session.add(new_log)
        db.session.commit()
        
    except Exception as e:
        print(f"Logging Error: {e}")

# --- STATS TRACKING (In-Memory for Dashboard Speed) ---
global_stats = {
    "text_gen": 0, "audio_gen": 0, "transcribe": 0, "pdf_gen": 0, 
    "chat_msgs": 0, "code_review": 0, "quiz_gen": 0,
    "file_conv": 0, "compression": 0, "vid_audio": 0
}

# --- CLEANUP TASK ---
def cleanup_old_files():
    try:
        now = time.time()
        # Check if directory exists before listing
        if os.path.exists(STATIC_FOLDER):
            for f in os.listdir(STATIC_FOLDER):
                fpath = os.path.join(STATIC_FOLDER, f)
                # Remove files older than 30 mins
                if os.path.isfile(fpath) and now - os.path.getmtime(fpath) > 1800:
                    try: os.remove(fpath)
                    except: pass
    except: pass

@app.before_request
def before_request_cleanup():
    # Only run cleanup on non-static requests to save resources
    if request.endpoint != 'serve_static_files': 
        threading.Thread(target=cleanup_old_files).start()

# --- AI HELPER FUNCTIONS ---
def clean_ai_text(text):
    if not text: return ""
    return text.replace("```html", "").replace("```json", "").replace("```", "").strip()

def get_groq_response(system_prompt, user_prompt, temperature=0.5):
    if not API_KEY: 
        return "Error: GROQ_API_KEY not found in .env file."
    try:
        client = Groq(api_key=API_KEY)
        completion = client.chat.completions.create(
            model="llama-3.3-70b-versatile",
            messages=[
                {"role": "system", "content": system_prompt}, 
                {"role": "user", "content": user_prompt}
            ],
            temperature=temperature, 
            max_tokens=2048
        )
        return completion.choices[0].message.content
    except Exception as e: 
        logging.error(f"AI Error: {e}")
        return f"AI Service Error: {str(e)}"

# ==============================================================================
#                               AUTH ROUTES
# ==============================================================================

@app.route('/')
def index():
    if not session.get('user_name') and not session.get('is_admin'):
        return redirect(url_for('login_page'))
    
    current_user = session.get('user_name', 'Guest')
    is_admin = session.get('is_admin', False)
    if is_admin: current_user = 'Administrator'
    
    return render_template('index.html', user_name=current_user, is_admin=is_admin)

@app.route('/login-page')
def login_page():
    if session.get('user_name') or session.get('is_admin'):
        return redirect(url_for('index'))
    return render_template('login.html')

@app.route('/login', methods=['POST'])
def login():
    data = request.json
    username = data.get('username')
    password = data.get('password')
    role = data.get('role', 'user')

    # Admin Login
    if role == 'admin':
        if username == ADMIN_USER and password == ADMIN_PASS:
            session['is_admin'] = True
            session['user_name'] = "Administrator"
            return jsonify({"success": True})
        return jsonify({"success": False, "error": "Invalid Admin Credentials"}), 401
    
    # User Login
    else:
        try:
            user = User.query.filter_by(username=username).first()
            if user and check_password_hash(user.password_hash, password):
                session['user_name'] = user.username
                session['is_admin'] = False
                
                # --- NEW: Log login activity ---
                log_activity('login', 'User Logged In')
                
                return jsonify({"success": True})
            return jsonify({"success": False, "error": "Invalid Username or Password"}), 401
        except Exception as e:
            # Re-create DB if it was wiped by Vercel
            db.create_all()
            return jsonify({"success": False, "error": "Database reset. Please Register again."}), 401

@app.route('/register', methods=['POST'])
def register():
    data = request.json
    username = data.get('username')
    password = data.get('password')

    if not username or not password:
        return jsonify({"success": False, "error": "Username and password required"}), 400
    
    # Ensure tables exist (Vercel fix)
    try:
        if User.query.filter_by(username=username).first():
            return jsonify({"success": False, "error": "Username already taken"}), 400
    except:
        db.create_all()
    
    hashed_pw = generate_password_hash(password)
    new_user = User(username=username, password_hash=hashed_pw)
    
    try:
        db.session.add(new_user)
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/reset-password', methods=['POST'])
def reset_password():
    data = request.json
    username = data.get('username')
    new_password = data.get('new_password')
    
    user = User.query.filter_by(username=username).first()
    if not user:
        return jsonify({"success": False, "error": "User not found"}), 404
        
    try:
        user.password_hash = generate_password_hash(new_password)
        db.session.commit()
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/logout', methods=['POST'])
def logout():
    session.clear()
    return jsonify({"success": True})

@app.route('/check-auth')
def check_auth():
    return jsonify({
        "is_admin": session.get('is_admin', False),
        "user_name": session.get('user_name', None)
    })

# ==============================================================================
#                               ADMIN & REPORTING
# ==============================================================================

# --- MODIFIED: Admin Dashboard Route ---
@app.route('/admin')
def admin_dashboard():
    # Admin Policy: Strict Access Control
    if not session.get('is_admin'):
        return redirect(url_for('login_page'))

    # We do NOT pass data here anymore. The frontend fetches it via AJAX.
    return render_template('admin.html', users=[], logs=[])

@app.route('/api/stats')
def get_stats():
    cpu, ram = 0, 0
    if session.get('is_admin', False):
        try: 
            cpu = psutil.cpu_percent(interval=None)
            ram = psutil.virtual_memory().percent
        except: pass
    return jsonify({"cpu": cpu, "ram": ram, "usage": global_stats})

# --- NEW: API to fetch Users for the Admin "Users" Tab ---
# In app.py, find the @app.route('/api/users') function
# Make sure it looks EXACTLY like this:

@app.route('/api/users')
def get_users_list():
    if not session.get('is_admin'): 
        return jsonify({"error": "Unauthorized"}), 403
    
    try:
        users = User.query.all()
        data = []
        for user in users:
            # Last Active Logic
            last_log = ActivityLog.query.filter_by(user_id=user.id).order_by(ActivityLog.timestamp.desc()).first()
            last_seen = "Never"
            if last_log:
                last_seen = last_log.timestamp.strftime("%Y-%m-%d %H:%M")
            
            # Action Count
            action_count = ActivityLog.query.filter_by(user_id=user.id).count()

            # --- THIS PART IS CRITICAL FOR FIXING "UNDEFINED" ---
            data.append({
                "id": user.id,
                "username": user.username,
                "full_name": user.full_name or "N/A",   # <--- Must be here
                "department": user.department or "N/A", # <--- Must be here
                "last_seen": last_seen,
                "action_count": action_count
            })
            # ----------------------------------------------------

        return jsonify(data)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

# --- NEW: API to fetch Activity Logs for the Admin "Activity" Tab ---
@app.route('/api/activity-logs')
def get_activity_logs_json():
    if not session.get('is_admin'): 
        return jsonify({"error": "Unauthorized"}), 403
    
    try:
        # Get last 100 logs
        logs = ActivityLog.query.order_by(ActivityLog.timestamp.desc()).limit(100).all()
        data = []
        for log in logs:
            username = log.user.username if log.user else "Admin/Guest"
            data.append({
                "date": log.timestamp.strftime("%Y-%m-%d %H:%M"),
                "user": username,
                "action": log.activity_type,
                "details": log.details
            })
        return jsonify(data)
    except Exception as e:
        return jsonify({"error": str(e)}), 500

@app.route('/download-report')
def download_report():
    if not session.get('is_admin'): return "Unauthorized", 401
    
    now = datetime.datetime.now().strftime("%Y-%m-%d %H:%M:%S")
    
    # Fetch real activity logs from DB
    try:
        logs = ActivityLog.query.order_by(ActivityLog.timestamp.desc()).limit(100).all()
    except:
        logs = []
    
    report_text = f"AI WORKSPACE - SYSTEM REPORT\nGenerated: {now}\n\n"
    
    report_text += "--- USAGE SUMMARY ---\n"
    for k, v in global_stats.items():
        report_text += f"{k.replace('_', ' ').title()}: {v}\n"
    
    report_text += "\n--- RECENT USER ACTIVITY (Last 100 Actions) ---\n"
    report_text += f"{'TIMESTAMP':<25} | {'USER':<20} | {'ACTION':<20} | {'DETAILS'}\n"
    report_text += "-" * 100 + "\n"
    
    for log in logs:
        username = log.user.username if log.user else "Admin/Guest"
        report_text += f"{str(log.timestamp):<25} | {username:<20} | {log.activity_type:<20} | {log.details}\n"
    
    return Response(
        report_text, 
        mimetype="text/plain", 
        headers={"Content-disposition": "attachment; filename=System_Report.txt"}
    )

# ==============================================================================
#                       VERCEL STATIC FILE SERVING (CRITICAL)
# ==============================================================================

# This route is REQUIRED to download files generated in the temp folder on Vercel
@app.route('/static/<path:filename>')
def serve_static_files(filename):
    return send_from_directory(STATIC_FOLDER, filename)

# ==============================================================================
#                               AI TOOLS
# ==============================================================================

@app.route('/chat', methods=['POST'])
def chat():
    log_activity('chat_msgs', 'User sent a chat message')
    msg = request.form.get('message', '')
    history = session.get('chat_history', [])
    sys_msg = {"role": "system", "content": "You are a helpful AI assistant."}
    try:
        messages = [sys_msg] + history + [{"role": "user", "content": msg}]
        client = Groq(api_key=API_KEY)
        completion = client.chat.completions.create(model="llama-3.3-70b-versatile", messages=messages, temperature=0.7)
        ai_reply = completion.choices[0].message.content
        history.append({"role": "user", "content": msg})
        history.append({"role": "assistant", "content": ai_reply})
        if len(history) > 6: history = history[-6:]
        session['chat_history'] = history
        return jsonify({"success": True, "response": ai_reply})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/clear-chat', methods=['POST'])
def clear_chat():
    session.pop('chat_history', None)
    return jsonify({"success": True})

@app.route('/generate-minutes', methods=['POST'])
def generate_minutes():
    log_activity('text_gen', 'Generated Meeting Minutes')
    notes = request.form.get('notes', '')
    res = get_groq_response("Convert notes to minutes.", notes)
    return jsonify({"success": True, "minutes": res})

@app.route('/generate-email', methods=['POST'])
def generate_email():
    log_activity('text_gen', 'Generated Email Draft')
    recipient = request.form.get('recipient')
    topic = request.form.get('topic')
    res = get_groq_response("Write a professional email.", f"To: {recipient} Topic: {topic}")
    return jsonify({"success": True, "email_content": res})

@app.route('/review-code', methods=['POST'])
def review_code():
    log_activity('code_review', 'Performed Code Review')
    code = request.form.get('code', '')
    res = get_groq_response("Review code.", code)
    return jsonify({"success": True, "review": res})

@app.route('/translate', methods=['POST'])
def translate():
    target = request.form.get('target_language')
    log_activity('text_gen', f'Translated text to {target}')
    text = request.form.get('text')
    res = get_groq_response(f"Translate to {target}.", text)
    return jsonify({"success": True, "translation": res})

@app.route('/generate-quiz', methods=['POST'])
def generate_quiz():
    log_activity('quiz_gen', 'Generated a Quiz')
    
    topic = request.form.get('topic', 'General Knowledge')
    count = request.form.get('count', '5')
    
    # 1. Strict HTML Prompt to AI
    prompt = (
        f"Create a {count}-question Multiple Choice Quiz about '{topic}'.\n"
        "Output ONLY raw HTML content (no ```html fences).\n"
        "Use this EXACT structure for every question:\n"
        "<div class='question-box'>\n"
        "  <h3 class='q-title'>1. Question text here?</h3>\n"
        "  <ul class='options-list'>\n"
        "    <li>A) Option 1</li>\n"
        "    <li>B) Option 2</li>\n"
        "    <li>C) Option 3</li>\n"
        "    <li>D) Option 4</li>\n"
        "  </ul>\n"
        "</div>\n"
        "At the very end, add: <h4>Answer Key</h4>\n"
        "<table class='answer-key'>...</table>"
    )
    
    raw_res = get_groq_response("You are a strict HTML quiz generator.", prompt)
    if not raw_res: return jsonify({"success": False, "error": "AI Failed"})
    
    clean_html = clean_ai_text(raw_res)
    
    # 2. PDF CSS Wrapper
    pdf_html = f"""
    <html>
    <head>
        <style>
            @page {{ size: A4; margin: 2cm; }}
            body {{ font-family: Helvetica, sans-serif; font-size: 12px; color: #000; line-height: 1.4; }}
            h1 {{ text-align: center; color: #4f46e5; border-bottom: 2px solid #eee; padding-bottom: 10px; margin-bottom: 20px; }}
            .question-box {{ margin-bottom: 10px; page-break-inside: avoid; }}
            .q-title {{ font-size: 14px; font-weight: bold; margin-bottom: 5px; color: #333; }}
            ul.options-list {{ margin: 0; padding-left: 20px; list-style-type: none; }}
            li {{ margin-bottom: 2px; padding: 2px 0; }}
            h4 {{ margin-top: 20px; border-bottom: 1px solid #ccc; }}
            table.answer-key {{ width: 100%; border-collapse: collapse; margin-top: 10px; font-size: 11px; }}
            th, td {{ border: 1px solid #999; padding: 5px; text-align: left; }}
            th {{ background-color: #f0f0f0; }}
        </style>
    </head>
    <body>
        <h1>Quiz: {topic}</h1>
        {clean_html}
    </body>
    </html>
    """
    
    fname = f"quiz_{uuid.uuid4().hex[:8]}.pdf"
    path = os.path.join(STATIC_FOLDER, fname)
    
    try:
        with open(path, "w+b") as f:
            pisa.CreatePDF(BytesIO(pdf_html.encode('utf-8')), dest=f)
        return jsonify({"success": True, "quiz": clean_html, "file_url": f"/static/{fname}"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/make-ppt', methods=['POST'])
def make_ppt():
    # 1. Get Data from Request
    topic = request.form.get('topic', 'Presentation')
    source_text = request.form.get('source_text', '') # <--- NOW USING YOUR INPUT TEXT

    log_activity('text_gen', f'Generated PPT: {topic}')
    
    # 2. Handle Template (Optional)
    template = request.files.get('template_file')
    prs = Presentation()
    if template:
        t_path = os.path.join(STATIC_FOLDER, f"temp_{uuid.uuid4()}.pptx")
        template.save(t_path)
        try: prs = Presentation(t_path)
        except: prs = Presentation()
        if os.path.exists(t_path): os.remove(t_path)
    
    # 3. Build Stronger AI Prompt
    # We combine the topic and the detailed source text
    content_input = f"TOPIC: {topic}\nDETAILS: {source_text}"
    
    system_instruction = (
        "You are a presentation generator. "
        "Convert the user's input into a slide deck structure. "
        "Strictly follow this format for every slide:\n"
        "SLIDE: [Title of the Slide]\n"
        "POINT: [Bullet point content]\n"
        "POINT: [Bullet point content]\n"
        "Do not output any conversational text, only the slide structure."
    )

    # 4. Get AI Response
    ai_text = get_groq_response(system_instruction, content_input)
    clean_response = clean_ai_text(ai_text)
    
    # 5. Parse and Build Slides
    slide = None
    for line in clean_response.split('\n'):
        line = line.strip()
        
        # Check for Slide Title (Case Insensitive)
        if line.upper().startswith("SLIDE:") or line.upper().startswith("SLIDE "):
            # Create a new slide (Layout 1 is usually Title + Content)
            layout = prs.slide_layouts[1] if len(prs.slide_layouts) > 1 else prs.slide_layouts[0]
            slide = prs.slides.add_slide(layout)
            
            # Extract title text
            title_text = line.split(':', 1)[-1].strip() if ':' in line else line
            try: slide.shapes.title.text = title_text
            except: pass
            
        # Check for Bullet Points (Handle 'POINT:', '-', or '*')
        elif (line.upper().startswith("POINT:") or line.startswith("-") or line.startswith("*")) and slide:
            # Extract point text
            if line.upper().startswith("POINT:"):
                point_text = line.split(':', 1)[-1].strip()
            else:
                point_text = line.lstrip("-* ").strip()
            
            try:
                # Add paragraph to the text body
                tf = slide.placeholders[1].text_frame
                p = tf.add_paragraph()
                p.text = point_text
                p.level = 0 # Top level bullet
            except: pass
            
    # 6. Save File
    fname = f"presentation_{uuid.uuid4().hex[:8]}.pptx"
    save_path = os.path.join(STATIC_FOLDER, fname)
    prs.save(save_path)
    
    return jsonify({"success": True, "file_url": f"/static/{fname}"})

@app.route('/text-to-audio', methods=['POST'])
def text_to_audio():
    log_activity('audio_gen', 'Converted Text to Speech')
    try:
        text = request.form.get('text')
        lang = request.form.get('target_language', 'en').split('-')[0]
        tts = gTTS(text=text, lang=lang)
        fname = f"audio_{uuid.uuid4().hex[:8]}.mp3"
        tts.save(os.path.join(STATIC_FOLDER, fname))
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/audio-to-text', methods=['POST'])
def audio_to_text():
    log_activity('transcribe', 'Transcribed Audio File')
    f = request.files['file']
    fname = f"temp_{uuid.uuid4().hex}.wav"
    path = os.path.join(STATIC_FOLDER, fname)
    f.save(path)
    try:
        r = sr.Recognizer()
        with sr.AudioFile(path) as src:
            txt = r.recognize_google(r.record(src), language=request.form.get('language', 'en-US'))
        os.remove(path)
        return jsonify({"success": True, "text": txt})
    except Exception as e: 
        if os.path.exists(path): os.remove(path)
        return jsonify({"success": False, "error": str(e)}), 500

@app.route('/convert-file', methods=['POST'])
def convert_file():
    fmt = request.form.get('format', 'PNG').upper()
    log_activity('file_conv', f'Converted file to {fmt}')
    f = request.files['file']
    try:
        img = PIL.Image.open(f)
        if fmt in ['JPG','JPEG']: img = img.convert('RGB')
        fname = f"conv_{uuid.uuid4().hex[:8]}.{fmt.lower()}"
        img.save(os.path.join(STATIC_FOLDER, fname), fmt)
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/compress-image', methods=['POST'])
def compress_image():
    log_activity('compression', 'Compressed Image')
    f = request.files['file']
    try:
        img = PIL.Image.open(f).convert('RGB')
        fname = f"comp_{uuid.uuid4().hex[:8]}.jpg"
        img.save(os.path.join(STATIC_FOLDER, fname), "JPEG", quality=30, optimize=True)
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/text-to-pdf', methods=['POST'])
def text_to_pdf():
    log_activity('pdf_gen', 'Generated PDF Document')
    fname = f"doc_{uuid.uuid4().hex[:8]}.pdf"
    try:
        with open(os.path.join(STATIC_FOLDER, fname), "w+b") as f:
            pisa.CreatePDF(BytesIO(request.form.get('html_content', '').encode('utf-8')), dest=f)
        return jsonify({"success": True, "file_url": f"/static/{fname}"})
    except Exception as e: return jsonify({"success": False, "error": str(e)}), 500

@app.route('/video-to-audio', methods=['POST'])
def video_to_audio():
    log_activity('vid_audio', 'Converted Video to Audio')
    
    if 'file' not in request.files: return jsonify({"success": False, "error": "No file"}), 400
    
    file = request.files['file']
    vid_name = f"temp_vid_{uuid.uuid4().hex[:8]}.mp4"
    vid_path = os.path.join(STATIC_FOLDER, vid_name)
    file.save(vid_path)
    
    audio_name = f"extracted_{uuid.uuid4().hex[:8]}.mp3"
    audio_path = os.path.join(STATIC_FOLDER, audio_name)
    
    try:
        # Added logger=None to prevent MoviePy from trying to write log files to read-only directories
        with VideoFileClip(vid_path) as clip:
            clip.audio.write_audiofile(audio_path, logger=None)
        return jsonify({"success": True, "file_url": f"/static/{audio_name}"})
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500
    finally:
        if os.path.exists(vid_path):
            try: os.remove(vid_path)
            except: pass

if __name__ == '__main__':
    app.run(host='0.0.0.0', port=5000, debug=True)